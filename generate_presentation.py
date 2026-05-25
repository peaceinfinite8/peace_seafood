import os
import sys
import subprocess

# 1. Pastikan python-pptx terinstall
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Modul 'python-pptx' belum terinstall. Menginstall otomatis...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Format Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # Definisikan Palet Warna Premium (Navy, Teal, Gold, White)
    COLOR_PRIMARY = RGBColor(15, 32, 67)     # Navy Gelap
    COLOR_ACCENT = RGBColor(6, 182, 212)    # Cyan/Teal
    COLOR_HIGHLIGHT = RGBColor(245, 158, 11) # Gold/Warning
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_DARK = RGBColor(33, 37, 41)
    COLOR_TEXT_MUTED = RGBColor(108, 117, 125)
    COLOR_SUCCESS = RGBColor(16, 185, 129)   # Hijau Sukses
    COLOR_DANGER = RGBColor(239, 68, 68)     # Merah Bahaya

    # Folder Screenshot
    assets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "presentation_assets")

    # Helper untuk membuat slide kosong (menghindari layout default PPT yang kaku)
    blank_layout = prs.slide_layouts[6]

    # Helper untuk menambahkan gambar secara aman jika ada
    def add_image_if_exists(slide, filename, left, top, width, height):
        full_path = os.path.join(assets_dir, filename) if assets_dir else filename
        if os.path.exists(full_path):
            slide.shapes.add_picture(full_path, left, top, width=width, height=height)
            print(f"Gambar {filename} berhasil dimasukkan ke slide.")
            return True
        else:
            # Jika gambar tidak ada, buat kotak placeholder abu-abu
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            shape.line.color.rgb = COLOR_TEXT_MUTED
            shape.line.width = Pt(1)
            
            txBox = slide.shapes.add_textbox(left, top + (height/2) - Inches(0.5), width, Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"[ TEMPELKAN SCREENSHOT DI SINI ]\n({filename})"
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT_MUTED
            print(f"Placeholder dibuat untuk {filename} (File tidak ditemukan).")
            return False

    # Helper untuk menambahkan header standar halaman konten
    def add_slide_header(slide, title_text, category_text="PEACE SEAFOOD"):
        # Garis aksen atas
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_ACCENT
        top_bar.line.fill.background()

        # Kategori Kecil
        cat_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.3))
        p_cat = cat_box.text_frame.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT

        # Judul Utama Slide
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
        p_title = title_box.text_frame.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY

    # ==========================================
    # SLIDE 1: COVER (DARK THEME)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_PRIMARY
    bg1.line.fill.background()

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "PROPOSAL DIGITALISASI OPERASIONAL"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_ACCENT
    
    p2 = tf1.add_paragraph()
    p2.text = "GUDANG IKAN & COLD STORAGE"
    p2.font.name = "Arial"
    p2.font.size = Pt(44)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "Mencegah Kebocoran Finansial, Meningkatkan Akurasi Timbangan, & Kontrol Real-Time"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(170, 190, 220)
    p3.space_before = Pt(20)

    footer_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.33), Inches(1.0))
    p_foot = footer_box.text_frame.paragraphs[0]
    p_foot.text = "Platform Sistem Terintegrasi: PEACE SEAFOOD\nDitawarkan Oleh: Project Manager & System Analyst Team"
    p_foot.font.name = "Arial"
    p_foot.font.size = Pt(12)
    p_foot.font.color.rgb = COLOR_WHITE

    # ==========================================
    # SLIDE 2: MASALAH OPERASIONAL GUDANG
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "Tantangan Klasik Gudang Ikan Tradisional", "PENGANTAR MASALAH")

    # Layout 3 Kolom Masalah
    problems = [
        ("❌ Kebocoran Timbangan", "Pencatatan berat manual di kertas papan jalan rawan salah tulis, terkena air, atau sengaja dimanipulasi oleh oknum (Fraud).", COLOR_DANGER),
        ("❌ Buta Kapasitas Gudang", "Tidak tahu pasti sisa kapasitas ruang beku (Cold Storage) tanpa harus masuk secara fisik ke dalam ruangan pendingin.", COLOR_PRIMARY),
        ("❌ Keuangan Terlambat", "Laporan laba/rugi baru bisa diketahui di akhir bulan secara manual. Nota piutang sering terselip dan telat ditagih.", COLOR_HIGHLIGHT)
    ]

    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    left_start = Inches(0.8)

    for i, (title, desc, color) in enumerate(problems):
        left = left_start + i * (col_width + col_gap)
        
        # Kotak Latar
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), col_width, Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 249, 250)
        card.line.color.rgb = RGBColor(222, 226, 230)
        card.line.width = Pt(1.5)

        # Border Aksen Atas Card
        accent_card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.2), col_width, Inches(0.15))
        accent_card.fill.solid()
        accent_card.fill.fore_color.rgb = color
        accent_card.line.fill.background()

        # Teks Judul & Deskripsi
        text_box = slide2.shapes.add_textbox(left + Inches(0.2), Inches(2.5), col_width - Inches(0.4), Inches(3.6))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = COLOR_TEXT_DARK
        p_desc.space_before = Pt(15)

    # ==========================================
    # SLIDE 3: SOLUSI - PEACE SEAFOOD
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "Digitalisasi Terintegrasi Bersama Peace Seafood", "SOLUSI DIGITAL")

    # Kiri: Deskripsi Solusi
    left_box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True

    p_sol = tf_left.paragraphs[0]
    p_sol.text = "Transformasi Total dari Kertas ke Sistem"
    p_sol.font.name = "Arial"
    p_sol.font.size = Pt(20)
    p_sol.font.bold = True
    p_sol.font.color.rgb = COLOR_PRIMARY

    bullets = [
        ("⚖️ Integrasi Timbangan Riil", "Penguncian data berat aktual langsung saat bongkar muat untuk mencegah manipulasi data."),
        ("❄️ Visibilitas Ruang Penyimpanan", "Pemantauan visual kapasitas Cold Storage secara instan lewat layar komputer atau ponsel."),
        ("📊 Keuangan Otomatis", "Arus kas, penjualan harian, sisa piutang, dan laba bersih langsung terhitung otomatis tanpa kalkulator manual."),
        ("🛡️ Akses Bertingkat Aman", "Pembagian wewenang yang jelas antara Bos, Admin Kantor, dan Checker Lapangan.")
    ]

    for b_title, b_desc in bullets:
        p_b = tf_left.add_paragraph()
        p_b.text = f"• {b_title}"
        p_b.font.name = "Arial"
        p_b.font.size = Pt(15)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_ACCENT
        p_b.space_before = Pt(12)

        p_bd = tf_left.add_paragraph()
        p_bd.text = b_desc
        p_bd.font.name = "Arial"
        p_bd.font.size = Pt(12)
        p_bd.font.color.rgb = COLOR_TEXT_MUTED
        p_bd.space_before = Pt(2)
        p_bd.level = 0

    # Kanan: Placeholder/Gambar Dashboard Utama
    add_image_if_exists(slide3, "01_dashboard_waterfall.png", Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 4: DASHBOARD EKSEKUTIF (LABA/RUGI)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "Pantau Laba Bersih Setiap Hari Tanpa Ribet", "MODUL BOS - KEUANGAN")

    # Kiri: Detail Laba/Rugi
    left_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_left4 = left_box4.text_frame
    tf_left4.word_wrap = True

    p_lr = tf_left4.paragraphs[0]
    p_lr.text = "Visualisasi Arus Kas Premium"
    p_lr.font.name = "Arial"
    p_lr.font.size = Pt(20)
    p_lr.font.bold = True
    p_lr.font.color.rgb = COLOR_PRIMARY

    p_lr_desc = tf_left4.add_paragraph()
    p_lr_desc.text = "Sistem menyediakan grafik khusus Waterfall Arus Kas untuk memudahkan Anda memahami perputaran uang di gudang secara visual:"
    p_lr_desc.font.name = "Arial"
    p_lr_desc.font.size = Pt(13)
    p_lr_desc.font.color.rgb = COLOR_TEXT_DARK
    p_lr_desc.space_before = Pt(10)

    lr_bullets = [
        ("Grafik Waterfall Interaktif", "Menunjukkan detail kenaikan kas dari penjualan & piutang, dikurangi pengeluaran operasional & stok untuk menghasilkan Laba Bersih."),
        ("Komposisi Kas Masuk vs Keluar", "Visualisasi pie/doughnut chart untuk memantau perbandingan kas secara harian."),
        ("Akses Eksklusif Bos", "Data keuangan bersifat privat dan hanya bisa diakses oleh Pemilik Gudang serta Super Admin.")
    ]

    for b_title, b_desc in lr_bullets:
        p_b = tf_left4.add_paragraph()
        p_b.text = f"✔ {b_title}"
        p_b.font.name = "Arial"
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_SUCCESS
        p_b.space_before = Pt(10)

        p_bd = tf_left4.add_paragraph()
        p_bd.text = b_desc
        p_bd.font.name = "Arial"
        p_bd.font.size = Pt(12)
        p_bd.font.color.rgb = COLOR_TEXT_MUTED
        p_bd.space_before = Pt(2)

    # Kanan: Tampilan Grafik Waterfall Laba/Rugi
    add_image_if_exists(slide4, "01_dashboard_waterfall.png", Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 5: MONITOR COLD STORAGE
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "Pantau Kapasitas Cold Storage Secara Visual", "MODUL GUDANG - STOK")

    # Kiri: Penjelasan Kapasitas
    left_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_left5 = left_box5.text_frame
    tf_left5.word_wrap = True

    p_cs = tf_left5.paragraphs[0]
    p_cs.text = "Keamanan Stok dan Kapasitas Ruang"
    p_cs.font.name = "Arial"
    p_cs.font.size = Pt(20)
    p_cs.font.bold = True
    p_cs.font.color.rgb = COLOR_PRIMARY

    cs_bullets = [
        ("Visual Gauge Real-Time", "Melihat sisa kapasitas penyimpanan beku langsung lewat diagram melengkung tanpa harus masuk ke dalam gudang."),
        ("Kapasitas Keamanan 5.000 kg", "Batas ukur penyimpanan yang disesuaikan secara otomatis berdasarkan data aktual di dalam database."),
        ("Peringatan Overcapacity", "Indikator otomatis berganti warna menjadi MERAH apabila penyimpanan sudah terisi di atas 80% guna menghindari kerusakan stok.")
    ]

    for b_title, b_desc in cs_bullets:
        p_b = tf_left5.add_paragraph()
        p_b.text = f"✔ {b_title}"
        p_b.font.name = "Arial"
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_ACCENT
        p_b.space_before = Pt(12)

        p_bd = tf_left5.add_paragraph()
        p_bd.text = b_desc
        p_bd.font.name = "Arial"
        p_bd.font.size = Pt(12)
        p_bd.font.color.rgb = COLOR_TEXT_MUTED
        p_bd.space_before = Pt(2)

    # Kanan: Grafik Gauge Kapasitas Gudang
    add_image_if_exists(slide5, "02_cold_storage_gauge.png", Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 6: ANTARMUKA WORKER GUDANG (MOBILE)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6, "Antarmuka Checker Lapangan yang Super Ringkas", "KEMUDAHAN OPERASIONAL")

    # Kiri: Penjelasan Checker
    left_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_left6 = left_box6.text_frame
    tf_left6.word_wrap = True

    p_ch = tf_left6.paragraphs[0]
    p_ch.text = "Sederhana, Cepat, dan Akurat"
    p_ch.font.name = "Arial"
    p_ch.font.size = Pt(20)
    p_ch.font.bold = True
    p_ch.font.color.rgb = COLOR_PRIMARY

    ch_bullets = [
        ("Didesain Khusus untuk HP/Tablet", "Checker lapangan dapat langsung menginput data berat timbangan langsung dari area bongkar muat ikan menggunakan ponsel."),
        ("Penguncian Berat Aktual (qty_actual)", "Mengurangi celah manipulasi data berat barang masuk karena data langsung terkirim secara real-time."),
        ("Minim Tombol & Teks Besar", "Meminimalisir kebingungan pekerja gudang beku yang kurang ramah terhadap teknologi.")
    ]

    for b_title, b_desc in ch_bullets:
        p_b = tf_left6.add_paragraph()
        p_b.text = f"✔ {b_title}"
        p_b.font.name = "Arial"
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_HIGHLIGHT
        p_b.space_before = Pt(12)

        p_bd = tf_left6.add_paragraph()
        p_bd.text = b_desc
        p_bd.font.name = "Arial"
        p_bd.font.size = Pt(12)
        p_bd.font.color.rgb = COLOR_TEXT_MUTED
        p_bd.space_before = Pt(2)

    # Kanan: Tampilan Mobile Input Timbangan
    add_image_if_exists(slide6, "03_checker_mobile.png", Inches(7.5), Inches(2.0), Inches(4.3), Inches(4.5))

    # ==========================================
    # SLIDE 7: NOTIFIKASI KRISIS & DETEKSI DINI
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide7, "Sistem Alarm: Cegah Kehabisan Stok & Piutang Macet", "KEAMANAN BISNIS")

    # Kiri: Deteksi Alarm
    left_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_left7 = left_box7.text_frame
    tf_left7.word_wrap = True

    p_al = tf_left7.paragraphs[0]
    p_al.text = "Sistem Proteksi Keuangan & Stok"
    p_al.font.name = "Arial"
    p_al.font.size = Pt(20)
    p_al.font.bold = True
    p_al.font.color.rgb = COLOR_PRIMARY

    al_bullets = [
        ("Peringatan Stok Menipis", "Notifikasi instan otomatis berkedip MERAH jika persediaan jenis ikan tertentu berada di bawah batas aman agar segera reorder."),
        ("Pengingat Piutang Jatuh Tempo", "Menghindari piutang pelanggan terselip. Sistem otomatis melacak tanggal tagihan dan memunculkan notifikasi jatuh tempo."),
        ("Mencegah Kerugian Finansial", "Menjaga perputaran kas gudang tetap sehat dan meminimalisir peluang hilangnya pelanggan karena stok kosong.")
    ]

    for b_title, b_desc in al_bullets:
        p_b = tf_left7.add_paragraph()
        p_b.text = f"✔ {b_title}"
        p_b.font.name = "Arial"
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_DANGER
        p_b.space_before = Pt(12)

        p_bd = tf_left7.add_paragraph()
        p_bd.text = b_desc
        p_bd.font.name = "Arial"
        p_bd.font.size = Pt(12)
        p_bd.font.color.rgb = COLOR_TEXT_MUTED
        p_bd.space_before = Pt(2)

    # Kanan: Peringatan Stok Menipis & Jatuh Tempo
    add_image_if_exists(slide7, "04_alert_krisis.png", Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5))

    # ==========================================
    # SLIDE 8: KENYAMANAN TEMA (LIGHT VS DARK MODE)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide8, "User Experience Fleksibel: Light & Dark Mode", "ESTETIKA & KENYAMANAN")

    # Penjelasan Ringkas Atas
    text_box8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0))
    tf8 = text_box8.text_frame
    tf8.word_wrap = True
    p_desc8 = tf8.paragraphs[0]
    p_desc8.text = "Operasional gudang beku sering kali berpencahayaan redup. Peace Seafood menyediakan fitur perubahan tema (Light/Dark Mode) yang otomatis merubah warna visual sistem secara instan demi kenyamanan mata dan visibilitas terbaik staf Anda di lapangan."
    p_desc8.font.name = "Arial"
    p_desc8.font.size = Pt(13)
    p_desc8.font.color.rgb = COLOR_TEXT_DARK

    # Tampilan Kiri: Light Mode
    left_img = slide8.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5.6), Inches(0.4))
    p_li = left_img.text_frame.paragraphs[0]
    p_li.text = "☀ LIGHT MODE (Tampilan Bersih & Terang)"
    p_li.font.name = "Arial"
    p_li.font.size = Pt(12)
    p_li.font.bold = True
    p_li.font.color.rgb = COLOR_PRIMARY
    p_li.alignment = PP_ALIGN.CENTER
    add_image_if_exists(slide8, "05_light_vs_dark.png", Inches(0.8), Inches(2.9), Inches(5.6), Inches(3.6))

    # Tampilan Kanan: Dark Mode
    right_img = slide8.shapes.add_textbox(Inches(6.8), Inches(2.5), Inches(5.6), Inches(0.4))
    p_ri = right_img.text_frame.paragraphs[0]
    p_ri.text = "🌙 DARK MODE (Nyaman di Lingkungan Redup Gudang)"
    p_ri.font.name = "Arial"
    p_ri.font.size = Pt(12)
    p_ri.font.bold = True
    p_ri.font.color.rgb = COLOR_ACCENT
    p_ri.alignment = PP_ALIGN.CENTER
    add_image_if_exists(slide8, "05_light_vs_dark.png", Inches(6.8), Inches(2.9), Inches(5.6), Inches(3.6))

    # ==========================================
    # SLIDE 9: RENCANA IMPLEMENTASI & PENUTUP (DARK THEME)
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    bg9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg9.fill.solid()
    bg9.fill.fore_color.rgb = COLOR_PRIMARY
    bg9.line.fill.background()

    # Kiri: Langkah Migrasi
    left_box9 = slide9.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(5.5), Inches(5.0))
    tf_left9 = left_box9.text_frame
    tf_left9.word_wrap = True

    p_cl = tf_left9.paragraphs[0]
    p_cl.text = "Mengapa Memilih Kami?"
    p_cl.font.name = "Arial"
    p_cl.font.size = Pt(24)
    p_cl.font.bold = True
    p_cl.font.color.rgb = COLOR_ACCENT

    bullets9 = [
        ("Transirasi Bertahap Aman", "Sistem kertas manual berjalan paralel bersama sistem digital selama 1-2 minggu masa penyesuaian agar operasional tidak terganggu."),
        ("Pelatihan Staf Lapangan", "Kami melatih petugas lapangan langsung di lokasi gudang hingga mereka mahir mencatat data di ponsel kurang dari 30 detik."),
        ("Migrasi Data Cepat", "Seluruh data supplier, produk ikan beku, dan pembeli langganan Anda akan diimpor langsung oleh tim kami.")
    ]

    for b_title, b_desc in bullets9:
        p_b = tf_left9.add_paragraph()
        p_b.text = f"✔ {b_title}"
        p_b.font.name = "Arial"
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_WHITE
        p_b.space_before = Pt(15)

        p_bd = tf_left9.add_paragraph()
        p_bd.text = b_desc
        p_bd.font.name = "Arial"
        p_bd.font.size = Pt(12)
        p_bd.font.color.rgb = RGBColor(170, 190, 220)
        p_bd.space_before = Pt(2)

    # Kanan: Ajakan Demo Penutup
    right_box9 = slide9.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.0), Inches(4.0))
    tf_right9 = right_box9.text_frame
    tf_right9.word_wrap = True

    p_r1 = tf_right9.paragraphs[0]
    p_r1.text = "SAATNYA MENGUNCI"
    p_r1.font.name = "Arial"
    p_r1.font.size = Pt(20)
    p_r1.font.bold = True
    p_r1.font.color.rgb = COLOR_HIGHLIGHT

    p_r2 = tf_right9.add_paragraph()
    p_r2.text = "KEBOCORAN FINANSIAL"
    p_r2.font.name = "Arial"
    p_r2.font.size = Pt(36)
    p_r2.font.bold = True
    p_r2.font.color.rgb = COLOR_WHITE
    p_r2.space_before = Pt(5)

    p_r3 = tf_right9.add_paragraph()
    p_r3.text = "DAN MEMBAWA GUDANG IKAN ANDA KE ERA DIGITAL MODERN."
    p_r3.font.name = "Arial"
    p_r3.font.size = Pt(16)
    p_r3.font.color.rgb = RGBColor(170, 190, 220)
    p_r3.space_before = Pt(10)

    p_r4 = tf_right9.add_paragraph()
    p_r4.text = "\nMari Jadwalkan Demo Langsung Hari Ini!"
    p_r4.font.name = "Arial"
    p_r4.font.size = Pt(16)
    p_r4.font.bold = True
    p_r4.font.color.rgb = COLOR_ACCENT
    p_r4.space_before = Pt(15)

    # Simpan File Presentasi
    output_filename = "c:\\xamppp\\htdocs\\peace_seafood\\PROPOSAL_PRESENTASI_PEACE_SEAFOOD.pptx"
    prs.save(output_filename)
    print(f"\n=========================================")
    print(f"SUKSES: File presentasi berhasil disimpan!")
    print(f"Lokasi file: {output_filename}")
    print(f"=========================================")

if __name__ == "__main__":
    create_presentation()
